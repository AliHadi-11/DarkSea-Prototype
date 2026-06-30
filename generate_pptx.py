from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette (Dark Sea Theme) ──────────────────────────────────────────
NAVY        = RGBColor(0x05, 0x14, 0x2E)   # slide background
TEAL        = RGBColor(0x00, 0xB4, 0xD8)   # accent / headings
CYAN_LIGHT  = RGBColor(0x90, 0xE0, 0xEF)   # sub-headings / icons
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_PANEL  = RGBColor(0x0A, 0x2A, 0x4A)   # content box fill
GOLD        = RGBColor(0xFF, 0xC3, 0x00)   # highlight numbers
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, colour=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour

def box(slide, left, top, width, height, fill=DARK_PANEL, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape

def accent_bar(slide, top=Inches(0.55), colour=TEAL):
    bar = slide.shapes.add_shape(1,
        Inches(0), top, SLIDE_W, Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = colour
    bar.line.fill.background()

def txbox(slide, text, left, top, width, height,
          font_size=18, bold=False, colour=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = colour
    return tb

def section_header(slide, title):
    """Teal left-bar + white title — used on content slides."""
    bar = slide.shapes.add_shape(1,
        Inches(0), Inches(0), Inches(0.25), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    txbox(slide, title, Inches(0.4), Inches(0.15),
          Inches(12.5), Inches(0.8),
          font_size=32, bold=True, colour=TEAL)
    accent_bar(slide, top=Inches(0.88))

def bullet_list(slide, items, left, top, width, height,
                font_size=16, icon="▪", colour=WHITE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = colour


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)

# Glowing title panel
box(sl, Inches(0.5), Inches(1.3), Inches(12.33), Inches(2.4), fill=DARK_PANEL, line=TEAL)
txbox(sl, "🌊  DarkSea", Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
      font_size=52, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)
txbox(sl, "A 3D First-Person Underwater Survival Horror Game",
      Inches(0.7), Inches(2.25), Inches(12), Inches(0.7),
      font_size=20, bold=False, colour=CYAN_LIGHT, align=PP_ALIGN.CENTER)

accent_bar(sl, top=Inches(3.8))

info = [
    ("Supervisor:", "TBD"),
    ("Group ID:",   "SP25-X"),
    ("Session:",    "2022–2026"),
]
y = Inches(4.0)
for label, val in info:
    txbox(sl, label, Inches(3.5), y, Inches(2), Inches(0.4),
          font_size=15, bold=True, colour=CYAN_LIGHT)
    txbox(sl, val,   Inches(5.5), y, Inches(4), Inches(0.4),
          font_size=15, colour=WHITE)
    y += Inches(0.45)

txbox(sl, "Student:  Ali  (Your Roll No)", Inches(3.5), y, Inches(6), Inches(0.4),
      font_size=15, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Problem Statement")

problem = (
    '"Millions of gamers seek immersive, story-driven survival experiences, yet the '
    'indie game market lacks titles that combine underwater horror, real-time oxygen '
    'management, and adaptive enemy AI within a single cohesive prototype. Existing '
    'games either focus purely on action or purely on exploration, leaving a gap for '
    'a tension-driven survival experience that tests both strategic thinking and '
    'reflexes. DarkSea bridges this gap by delivering a fully playable 3-level '
    'prototype with escalating AI behaviours, a resource-management system, and an '
    'immersive underwater atmosphere — built entirely in Unity."'
)
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.5), fill=DARK_PANEL, line=TEAL)
txbox(sl, problem, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.2),
      font_size=18, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Key Objectives
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Key Objectives")

objectives = [
    "Design and implement a 3-level 3D underwater survival game using Unity and C#.",
    "Develop an adaptive Enemy AI system with three distinct behavioural modes: Chase, Passive, and Hunt.",
    "Implement a real-time Oxygen Management system that acts as the primary survival pressure mechanic.",
    "Build a complete user authentication flow (Register / Login / Forgot Password) using PlayerPrefs.",
    "Create a progressive scoring system rewarding speed, oxygen conservation, and collectible retrieval.",
    "Produce a polished, fully playable prototype ready for FYP demonstration.",
]
y = Inches(1.05)
for i, obj in enumerate(objectives, 1):
    box(sl, Inches(0.4), y, Inches(12.5), Inches(0.72), fill=DARK_PANEL, line=TEAL)
    txbox(sl, f"0{i}", Inches(0.5), y + Inches(0.05), Inches(0.6), Inches(0.6),
          font_size=20, bold=True, colour=GOLD)
    txbox(sl, obj, Inches(1.15), y + Inches(0.08), Inches(11.6), Inches(0.6),
          font_size=15, colour=WHITE)
    y += Inches(0.84)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Motivation
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Motivation")

points = [
    ("Game Dev as a Career Path",
     "Unity is one of the most in-demand engines worldwide. This project builds marketable skills in 3D game development, AI programming, and UX design."),
    ("Under-explored Genre",
     "Underwater horror survival is a niche genre with high commercial potential. DarkSea explores this space as a technical and creative challenge."),
    ("AI Research Application",
     "Implementing NavMesh-based enemy AI with multiple behavioural states applies real AI/ML concepts in a tangible, interactive medium."),
    ("Full-Stack Game Architecture",
     "From auth systems to score persistence to scene management — this project demonstrates end-to-end software architecture skills."),
]
y = Inches(1.05)
for title, desc in points:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.2), fill=DARK_PANEL, line=TEAL)
    txbox(sl, f"🎯  {title}", Inches(0.6), y + Inches(0.05),
          Inches(12), Inches(0.45), font_size=16, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.48),
          Inches(12), Inches(0.6), font_size=14, colour=WHITE)
    y += Inches(1.38)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Game Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Game Overview")

overview_points = [
    "Genre: 3D First-Person Survival Horror (Underwater)",
    "Engine: Unity (C#)",
    "Perspective: First-Person with Rigidbody-based movement",
    "Levels: 3 fully designed levels with unique objectives",
    "Core Mechanic: Oxygen depletion — run out and it's Game Over",
    "Enemy AI: NavMeshAgent with Chase / Passive / Hunt modes",
    "Win Condition: Reach the level exit before oxygen runs out",
    "Progression: Unlock next level by completing current one; best scores saved locally",
]
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)
bullet_list(sl, overview_points, Inches(0.6), Inches(1.15),
            Inches(12.1), Inches(5.6), font_size=17, icon="🔹")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Target Audience
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Target Audience")

audience = [
    ("🎮  Casual Gamers",
     "Players who enjoy short-session, tension-filled games with clear win/lose conditions and simple controls."),
    ("👾  Horror / Survival Enthusiasts",
     "Fans of underwater horror and resource-management games such as Subnautica and Alien: Isolation."),
    ("🧑‍💻  Game Dev Students",
     "Computer Science students studying Unity, AI, or game design who can learn from the open prototype architecture."),
    ("🏫  Academic Evaluators",
     "FYP committee members assessing software engineering, AI implementation, and system design quality."),
]
y = Inches(1.05)
for icon_title, desc in audience:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.25), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_title, Inches(0.6), y + Inches(0.05),
          Inches(12), Inches(0.45), font_size=16, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.48),
          Inches(12), Inches(0.65), font_size=14, colour=WHITE)
    y += Inches(1.4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Proposed Solution / Key Features
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Proposed Solution")

features = [
    ("🌊 Oxygen System",
     "Real-time oxygen drain forces time-pressure decisions. Pickups and tank collection replenish oxygen mid-level."),
    ("🤖 Adaptive Enemy AI",
     "Three AI modes (Chase / Passive / Hunt) give each level a distinct feel. Hunt mode uses mimic voice to disorient the player."),
    ("📡 Sonar Mechanic",
     "SPACEBAR activates sonar — reveals wall distances and nearby enemies, but also risks triggering enemy alert."),
    ("🏆 Scoring & Persistence",
     "Score = oxygen remaining × 10 + time bonus + tanks collected × 50. Best scores saved via PlayerPrefs per level."),
]
y = Inches(1.05)
for icon_title, desc in features:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.3), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_title, Inches(0.6), y + Inches(0.05),
          Inches(12), Inches(0.45), font_size=16, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.5),
          Inches(12), Inches(0.65), font_size=14, colour=WHITE)
    y += Inches(1.45)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Gameplay Mechanics
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Gameplay Mechanics")

left_items = [
    "W / S — Move Forward / Backward",
    "A / D + Mouse — Rotate / Turn",
    "SPACEBAR — Activate Sonar pulse",
    "SPACEBAR (close range) — Kill enemy (≤ 3.5 m)",
    "ESC — Pause / Resume",
]
right_items = [
    "Oxygen drains continuously — warning at ≤ 20%",
    "Touch oxygen pickup → +40% oxygen refill",
    "Collect 3 tanks (Level 2) to unlock exit gate",
    "Reach exit trigger → level complete",
    "Enemy catch range: 2.5 m → instant Game Over",
]

box(sl, Inches(0.4),  Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)
box(sl, Inches(6.75), Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl, "🕹️  Controls", Inches(0.55), Inches(1.1),
      Inches(5.8), Inches(0.5), font_size=17, bold=True, colour=TEAL)
bullet_list(sl, left_items, Inches(0.55), Inches(1.65),
            Inches(5.8), Inches(4.8), font_size=15, icon="▸")

txbox(sl, "⚙️  Rules & Interactions", Inches(6.9), Inches(1.1),
      Inches(5.8), Inches(0.5), font_size=17, bold=True, colour=TEAL)
bullet_list(sl, right_items, Inches(6.9), Inches(1.65),
            Inches(5.8), Inches(4.8), font_size=15, icon="▸")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Scene / Game Flow
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Scene Flow")

scenes = [
    ("1", "RegisterScene",      "User creates account (username / email / password)"),
    ("2", "LoginScene",         "Login or recover password via email lookup"),
    ("3", "MainMenu",           "Welcome screen — Start Mission or Quit"),
    ("4", "Level1_Transition",  "3-second cutscene → auto-loads Level 1"),
    ("5", "DarkSea (Level 1)",  "Chase enemy, reach exit — earn score"),
    ("6", "Level2_Transition",  "Cutscene bridge → Level 2"),
    ("7", "Level_2",            "Passive enemy — collect 3 tanks to unlock exit"),
    ("8", "Level3_Transition",  "Cutscene bridge → Level 3"),
    ("9", "Level_3",            "Hunt enemy (mimic voice, faster) — reach exit"),
]
y = Inches(1.05)
for num, scene, desc in scenes:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(0.54), fill=DARK_PANEL, line=TEAL)
    txbox(sl, num,   Inches(0.5),  y + Inches(0.05), Inches(0.4), Inches(0.45),
          font_size=14, bold=True, colour=GOLD)
    txbox(sl, scene, Inches(0.95), y + Inches(0.05), Inches(3.2), Inches(0.45),
          font_size=14, bold=True, colour=TEAL)
    txbox(sl, desc,  Inches(4.2),  y + Inches(0.05), Inches(8.5), Inches(0.45),
          font_size=14, colour=WHITE)
    y += Inches(0.61)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — Oxygen System (Core Mechanic)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Core Mechanic — Oxygen System")

left = [
    "Oxygen drains at oxygenDecreaseSpeed units/second",
    "UI warning activates when oxygen ≤ 20%",
    "Oxygen hits 0 → ShowGameOver('OUT OF OXYGEN')",
    "OxygenPickup triggers RefillOxygen(+40f)",
    "CollectibleTank also refills oxygen on pickup",
    "Enemy catch also calls ShowGameOver(reason)",
    "Game Over sets Time.timeScale = 0 (freeze)",
    "RestartGame() reloads current scene and resets time",
]
right_stats = [
    ("Drain Rate",    "Configurable per level"),
    ("Warning Threshold", "≤ 20%"),
    ("Pickup Refill", "+40 units"),
    ("Script",        "OxygenSystem.cs"),
    ("Dependencies",  "EnemyAI_Final.cs, OxygenPickup.cs, CollectibleTank.cs"),
]

box(sl, Inches(0.4),  Inches(1.05), Inches(7.6), Inches(5.8), fill=DARK_PANEL, line=TEAL)
box(sl, Inches(8.2), Inches(1.05), Inches(4.7), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl, "How It Works", Inches(0.55), Inches(1.1), Inches(7.3), Inches(0.5),
      font_size=17, bold=True, colour=TEAL)
bullet_list(sl, left, Inches(0.55), Inches(1.65), Inches(7.3), Inches(5.0),
            font_size=14, icon="▸")

txbox(sl, "Quick Stats", Inches(8.35), Inches(1.1), Inches(4.4), Inches(0.5),
      font_size=17, bold=True, colour=TEAL)
ys = Inches(1.65)
for label, val in right_stats:
    txbox(sl, label, Inches(8.35), ys, Inches(4.4), Inches(0.3),
          font_size=13, bold=True, colour=CYAN_LIGHT)
    txbox(sl, val,   Inches(8.35), ys + Inches(0.3), Inches(4.4), Inches(0.35),
          font_size=13, colour=WHITE)
    ys += Inches(0.82)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — Enemy AI System
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Enemy AI System")

modes = [
    ("🔵  Passive Mode", "Patrol Mode",
     "Enemy patrols fixed waypoints. Never attacks the player. Used in Level 2 to create tension without direct danger."),
    ("🔴  Chase Mode", "Pursuit Mode",
     "Enemy detects player and pursues via NavMeshAgent pathfinding. Kills player if within catchRange (2.5 m). Used in Level 1."),
    ("🟠  Hunt Mode", "Advanced Pursuit",
     "Faster speed + periodic mimic voice audio clips to confuse/disorient. Used in Level 3 as the hardest enemy variant."),
]
y = Inches(1.05)
for icon_mode, subtitle, desc in modes:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.5), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_mode, Inches(0.6), y + Inches(0.05), Inches(6), Inches(0.5),
          font_size=17, bold=True, colour=TEAL)
    txbox(sl, subtitle, Inches(0.6), y + Inches(0.5), Inches(6), Inches(0.35),
          font_size=13, bold=True, colour=CYAN_LIGHT)
    txbox(sl, desc, Inches(0.6), y + Inches(0.85), Inches(12.1), Inches(0.55),
          font_size=14, colour=WHITE)
    y += Inches(1.65)

txbox(sl, "Script: EnemyAI_Final.cs  |  Navigation: Unity NavMeshAgent  |  Kill Range: 2.5 m",
      Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
      font_size=13, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — Level 1 Design
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Level Design — Level 1: DarkSea")

details = [
    ("Scene Name",    "DarkSea"),
    ("Objective",     "Navigate the underwater maze and reach the exit gate"),
    ("Enemy Type",    "Chase Mode — actively pursues the player"),
    ("AI Navigation", "NavMesh baked for DarkSea scene"),
    ("Win Condition", "Trigger LevelExit collider"),
    ("Lose Condition","Oxygen depleted OR enemy catches player"),
    ("Oxygen Pickups","Present throughout maze"),
    ("Sonar Use",     "SPACEBAR reveals walls + detects nearby enemy"),
    ("Score Formula", "Oxygen × 10 + Time Bonus"),
    ("Transition",    "Leads to Level2_Transition scene"),
]
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)
y = Inches(1.15)
for label, val in details:
    txbox(sl, f"{label}:", Inches(0.6), y, Inches(3.0), Inches(0.45),
          font_size=15, bold=True, colour=TEAL)
    txbox(sl, val, Inches(3.65), y, Inches(9.0), Inches(0.45),
          font_size=15, colour=WHITE)
    y += Inches(0.5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — Level 2 Design
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Level Design — Level 2: Collect & Survive")

details = [
    ("Scene Name",     "Level_2"),
    ("Objective",      "Collect 3 oxygen tanks scattered in the environment"),
    ("Enemy Type",     "Passive Mode — patrols but never attacks"),
    ("Mission Item",   "CollectibleTank — 3 required to unlock exit"),
    ("Win Condition",  "Collect all 3 tanks THEN trigger LevelExit"),
    ("Lose Condition", "Oxygen depleted"),
    ("Oxygen Refill",  "Each tank pickup also refills oxygen"),
    ("UI Counter",     "Tank counter updated on each pickup (TankCollector)"),
    ("Score Formula",  "Oxygen × 10 + Time Bonus + Tanks × 50"),
    ("Transition",     "Leads to Level3_Transition scene"),
]
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)
y = Inches(1.15)
for label, val in details:
    txbox(sl, f"{label}:", Inches(0.6), y, Inches(3.0), Inches(0.45),
          font_size=15, bold=True, colour=TEAL)
    txbox(sl, val, Inches(3.65), y, Inches(9.0), Inches(0.45),
          font_size=15, colour=WHITE)
    y += Inches(0.5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — Level 3 Design
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Level Design — Level 3: Hunt Mode")

details = [
    ("Scene Name",    "Level_3"),
    ("Objective",     "Survive the Hunt enemy and reach the exit"),
    ("Enemy Type",    "Hunt Mode — faster speed + mimic voice audio"),
    ("Mimic Voice",   "Enemy plays periodic audio clips to disorient player"),
    ("Win Condition", "Trigger LevelExit collider"),
    ("Lose Condition","Oxygen depleted OR enemy catches player"),
    ("Difficulty",    "Highest — fastest enemy, psychological pressure"),
    ("Oxygen Pickups","Present but harder to safely reach"),
    ("Score Formula", "Oxygen × 10 + Time Bonus"),
    ("Final Level",   "Completion shows Game Complete / Thank You screen"),
]
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)
y = Inches(1.15)
for label, val in details:
    txbox(sl, f"{label}:", Inches(0.6), y, Inches(3.0), Inches(0.45),
          font_size=15, bold=True, colour=TEAL)
    txbox(sl, val, Inches(3.65), y, Inches(9.0), Inches(0.45),
          font_size=15, colour=WHITE)
    y += Inches(0.5)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — Sonar & Player Controls Deep Dive
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Sonar System & Player Controller")

left = [
    "Script: PlayerMovementFinal.cs",
    "Movement: Rigidbody.MovePosition (W/S axis)",
    "Rotation: A/D + Mouse X (combined rotation)",
    "SPACEBAR — dual-function key:",
    "     • Sonar pulse: measures wall distances, detects enemy",
    "     • Kill: if enemy within killRange (3.5 m) → enemy destroyed",
    "[RequireComponent(Rigidbody)] enforced",
    "SonarSystem.cs exists but is superseded by PlayerMovementFinal",
]
right = [
    "Sonar outputs distance readout to screen UI",
    "Sonar detects EnemyAI_Final tagged objects",
    "KillRange check: Physics.OverlapSphere(3.5 m)",
    "Kill → LevelManager.EnemyKilled() → Win panel",
    "ESC → PauseManager toggles pause",
    "Pause: Time.timeScale = 0 / Resume: = 1",
    "GameMessage.Show() for on-screen flash alerts",
    "WaitForSecondsRealtime — works during pause",
]

box(sl, Inches(0.4),  Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)
box(sl, Inches(6.75), Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl, "Player Controller", Inches(0.55), Inches(1.1), Inches(5.8), Inches(0.5),
      font_size=16, bold=True, colour=TEAL)
bullet_list(sl, left, Inches(0.55), Inches(1.65), Inches(5.8), Inches(5.0),
            font_size=13, icon="▸")

txbox(sl, "Sonar Details", Inches(6.9), Inches(1.1), Inches(5.8), Inches(0.5),
      font_size=16, bold=True, colour=TEAL)
bullet_list(sl, right, Inches(6.9), Inches(1.65), Inches(5.8), Inches(5.0),
            font_size=13, icon="▸")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — Scoring System
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Scoring System")

formula_items = [
    ("Oxygen Bonus",  "Oxygen Remaining × 10",       "Rewards efficient oxygen management"),
    ("Time Bonus",    "Configurable per level",       "Rewards completing level quickly"),
    ("Tank Bonus",    "Tanks Collected × 50",         "Level 2 only — rewards full collection"),
    ("Best Score",    "Saved to PlayerPrefs",         'Key: "BestScore_L{levelNumber}"'),
    ("Level Unlock",  "PlayerPrefs[UnlockedLevel]",   "Increments on level completion"),
]

box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl, "Score  =  (Oxygen × 10)  +  Time Bonus  +  (Tanks × 50)",
      Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.65),
      font_size=20, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

y = Inches(1.9)
for comp, formula, note in formula_items:
    box(sl, Inches(0.6), y, Inches(12.1), Inches(0.75), fill=NAVY, line=TEAL)
    txbox(sl, comp,    Inches(0.7), y + Inches(0.05), Inches(2.8), Inches(0.65),
          font_size=14, bold=True, colour=TEAL)
    txbox(sl, formula, Inches(3.6), y + Inches(0.05), Inches(3.5), Inches(0.65),
          font_size=14, bold=True, colour=GOLD)
    txbox(sl, note,    Inches(7.2), y + Inches(0.05), Inches(5.3), Inches(0.65),
          font_size=13, colour=WHITE)
    y += Inches(0.9)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — Technology Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Technology Stack")

tech = [
    ("🎮  Game Engine",    "Unity (LTS)",              "Primary development environment; 3D scene management, physics, rendering"),
    ("💻  Language",       "C#",                       "All gameplay scripts written in C# (.cs) under Assets/GameScripts/"),
    ("🧭  Navigation",     "Unity NavMesh",            "NavMeshAgent drives EnemyAI_Final pathfinding across baked NavMesh"),
    ("💾  Persistence",    "Unity PlayerPrefs",        "Stores user accounts (plaintext prototype), scores, and level unlock state"),
    ("📦  Plugins",        "NuGet Packages",           "NuGet folder present under Assets/Plugins/ for extended library support"),
    ("🎨  Assets",         "Custom Textures/Materials","medieval_blocks_03_diff_4k.mat; brown mud 4K blend textures"),
]
y = Inches(1.05)
for icon_cat, tool, desc in tech:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(0.9), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_cat, Inches(0.55), y + Inches(0.05), Inches(2.5), Inches(0.4),
          font_size=14, bold=True, colour=CYAN_LIGHT)
    txbox(sl, tool,     Inches(3.1), y + Inches(0.05), Inches(2.5), Inches(0.4),
          font_size=14, bold=True, colour=GOLD)
    txbox(sl, desc,     Inches(5.7), y + Inches(0.05), Inches(7.0), Inches(0.8),
          font_size=13, colour=WHITE)
    y += Inches(1.03)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "System Architecture")

# Draw simple architecture boxes
layers = [
    (Inches(0.5),  Inches(1.2),  Inches(12.0), Inches(0.9), "UI Layer",
     "MainMenuController  |  PauseManager  |  OxygenSystem UI  |  GameMessage", TEAL),
    (Inches(0.5),  Inches(2.3),  Inches(12.0), Inches(0.9), "Gameplay Layer",
     "PlayerMovementFinal  |  OxygenSystem  |  EnemyAI_Final  |  LevelExit  |  LevelManager", TEAL),
    (Inches(0.5),  Inches(3.4),  Inches(12.0), Inches(0.9), "Collectibles & Items",
     "OxygenPickup  |  CollectibleTank  |  TankCollector (static counter)", TEAL),
    (Inches(0.5),  Inches(4.5),  Inches(12.0), Inches(0.9), "Scene Management",
     "TransitionManager  |  SceneManager.LoadScene()", TEAL),
    (Inches(0.5),  Inches(5.6),  Inches(12.0), Inches(0.9), "Auth & Persistence",
     "RegisterManager  |  LoginTabSystem  |  PlayerPrefs (local storage only)", TEAL),
]
for lft, top, w, h, title, content, col in layers:
    box(sl, lft, top, w, h, fill=DARK_PANEL, line=col)
    txbox(sl, title,   lft + Inches(0.1), top + Inches(0.04), Inches(2.8), Inches(0.4),
          font_size=14, bold=True, colour=col)
    txbox(sl, content, lft + Inches(3.0), top + Inches(0.04), Inches(8.8), Inches(0.8),
          font_size=13, colour=WHITE)
    # arrow down (except last)
    if top < Inches(5.6):
        arr = sl.shapes.add_shape(1,
            Inches(6.3), top + h, Inches(0.3), Inches(0.28))
        arr.fill.solid(); arr.fill.fore_color.rgb = TEAL
        arr.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — Authentication System
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Authentication System")

auth_points = [
    ("RegisterManager.cs",
     "Accepts: Username, Email, Password, Confirm Password\n"
     "Saves: PlayerPrefs[username] = password  |  PlayerPrefs[username_email] = email\n"
     "Redirects to: LoginScene after successful registration"),
    ("LoginTabSystem.cs",
     "Tab/Enter key support for smooth field navigation\n"
     "Forgot Password popup: shows password in plaintext if email matches stored value\n"
     "Sets: PlayerPrefs[CurrentPlayer] = username  →  used by MainMenuController"),
    ("MainMenuController.cs",
     'Displays: "Welcome, [CurrentPlayer]!" greeting text\n'
     "Buttons: Start Mission → Level1_Transition  |  Quit → Application.Quit()"),
]
y = Inches(1.05)
for script, desc in auth_points:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.55), fill=DARK_PANEL, line=TEAL)
    txbox(sl, script, Inches(0.6), y + Inches(0.05), Inches(12.1), Inches(0.45),
          font_size=16, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.5), Inches(12.1), Inches(1.0),
          font_size=14, colour=WHITE)
    y += Inches(1.72)

txbox(sl, "⚠️  Note: Passwords stored in plaintext via PlayerPrefs — intentional simplification for prototype scope.",
      Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
      font_size=12, colour=GOLD, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — Functional Requirements
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Functional Requirements")

reqs = [
    "Player can Register, Login, and recover password",
    "Player can move (W/S) and rotate (A/D + Mouse) in 3D space",
    "Oxygen depletes over time; game ends when oxygen reaches zero",
    "Player can collect oxygen pickups to replenish oxygen",
    "Enemy AI can patrol (Passive), chase (Chase), or hunt with audio (Hunt)",
    "Sonar (SPACEBAR) detects walls and nearby enemies",
    "Player can kill enemy via SPACEBAR within 3.5 m range",
    "Level 2: Player must collect 3 tanks before exit unlocks",
    "Level exit calculates and saves best score per level to PlayerPrefs",
    "Pause menu (ESC) provides Resume, Restart, and Main Menu options",
    "Game Over screen appears on oxygen depletion or enemy catch",
    "Level completion unlocks next level in PlayerPrefs",
]
box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)
bullet_list(sl, reqs, Inches(0.6), Inches(1.15),
            Inches(12.1), Inches(5.6), font_size=15, icon="✅")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — Non-Functional Requirements
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Non-Functional Requirements")

nfr = [
    ("⚡ Performance",
     "Game must maintain stable frame rate on mid-range PC hardware. Physics and NavMesh calculations must not cause perceptible lag."),
    ("🎮 Usability",
     "Controls must be intuitive. HUD (oxygen bar, tank counter, flash messages) must be readable at all times during gameplay."),
    ("🔒 Security (Prototype Scope)",
     "Auth is local PlayerPrefs — adequate for a prototype. No network calls, no server-side exposure."),
    ("🔧 Maintainability",
     "Scripts follow single-responsibility where possible. OxygenSystem, EnemyAI, and LevelExit are modular and independently editable."),
    ("📈 Scalability",
     "TransitionManager and LevelExit architecture support adding new levels by simply creating new scenes and updating nextSceneName."),
]
y = Inches(1.05)
for icon_title, desc in nfr:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(1.1), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_title, Inches(0.6), y + Inches(0.05), Inches(12.1), Inches(0.4),
          font_size=15, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.48), Inches(12.1), Inches(0.55),
          font_size=13, colour=WHITE)
    y += Inches(1.25)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — System Requirements
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "System Requirements")

hw = [
    "CPU: Dual-core 2.0 GHz or higher",
    "RAM: 4 GB minimum (8 GB recommended)",
    "GPU: DirectX 11 compatible graphics card",
    "Storage: ~500 MB for project build",
    "Input: Keyboard + Mouse",
    "Display: 1280 × 720 minimum resolution",
]
sw = [
    "OS: Windows 10 / 11 (64-bit)",
    "Unity: 2021 LTS or newer (for development)",
    "Runtime: Unity Player (bundled in build)",
    ".NET: .NET Standard 2.1 (via Unity)",
    "NavMesh: Baked into DarkSea scene",
    "No internet connection required (offline game)",
]

box(sl, Inches(0.4),  Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)
box(sl, Inches(6.75), Inches(1.05), Inches(6.1), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl, "🖥️  Hardware Requirements", Inches(0.55), Inches(1.1), Inches(5.8), Inches(0.5),
      font_size=16, bold=True, colour=TEAL)
bullet_list(sl, hw, Inches(0.55), Inches(1.65), Inches(5.8), Inches(5.0),
            font_size=15, icon="▪")

txbox(sl, "💻  Software Requirements", Inches(6.9), Inches(1.1), Inches(5.8), Inches(0.5),
      font_size=16, bold=True, colour=TEAL)
bullet_list(sl, sw, Inches(6.9), Inches(1.65), Inches(5.8), Inches(5.0),
            font_size=15, icon="▪")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 23 — Use Case Diagram (placeholder)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Use Case Diagram")

box(sl, Inches(0.4), Inches(1.05), Inches(12.5), Inches(5.8), fill=DARK_PANEL, line=TEAL)

txbox(sl,
      "A use case diagram shows how different actors interact with the DarkSea system.\n\n"
      "Actors:\n"
      "  • Player (primary) — interacts with all gameplay systems\n"
      "  • System / Unity Engine — manages scene transitions, physics, NavMesh\n\n"
      "Key Use Cases:\n"
      "  • Register / Login / Recover Password\n"
      "  • Start Game / Navigate Level / Collect Items\n"
      "  • Activate Sonar / Kill Enemy\n"
      "  • Pause / Resume / Restart Game\n"
      "  • Complete Level / View Score / Unlock Next Level\n\n"
      "[Insert formal UML Use Case Diagram here]",
      Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.5),
      font_size=16, colour=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 24 — SWOT Analysis
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "SWOT Analysis")

quad_col = RGBColor(0x08, 0x20, 0x40)

# Top-left: Strengths
box(sl, Inches(0.4),  Inches(1.05), Inches(6.1), Inches(2.8), fill=quad_col, line=TEAL)
txbox(sl, "💪  Strengths", Inches(0.55), Inches(1.1), Inches(5.8), Inches(0.45),
      font_size=15, bold=True, colour=TEAL)
bullet_list(sl,
    ["Original game concept with unique mechanics",
     "Full 3-level playable prototype",
     "Modular, extensible C# codebase",
     "Adaptive enemy AI with 3 modes"],
    Inches(0.55), Inches(1.55), Inches(5.8), Inches(2.1), font_size=13, icon="▪")

# Top-right: Weaknesses
box(sl, Inches(6.75), Inches(1.05), Inches(6.1), Inches(2.8), fill=quad_col, line=GOLD)
txbox(sl, "⚠️  Weaknesses", Inches(6.9), Inches(1.1), Inches(5.8), Inches(0.45),
      font_size=15, bold=True, colour=GOLD)
bullet_list(sl,
    ["PlayerPrefs auth — not production-grade",
     "Sonar logic duplicated across two scripts",
     "No sound/music system beyond mimic voice",
     "Limited level variety (3 levels)"],
    Inches(6.9), Inches(1.55), Inches(5.8), Inches(2.1), font_size=13, icon="▪")

# Bottom-left: Opportunities
box(sl, Inches(0.4),  Inches(4.05), Inches(6.1), Inches(2.8), fill=quad_col, line=RGBColor(0x00,0xC8,0x55))
txbox(sl, "🌟  Opportunities", Inches(0.55), Inches(4.1), Inches(5.8), Inches(0.45),
      font_size=15, bold=True, colour=RGBColor(0x00,0xC8,0x55))
bullet_list(sl,
    ["Expand to 5+ levels with new mechanics",
     "Add multiplayer / co-op survival mode",
     "Publish as indie game on itch.io / Steam",
     "Integrate backend auth (Firebase / Supabase)"],
    Inches(0.55), Inches(4.55), Inches(5.8), Inches(2.1), font_size=13, icon="▪")

# Bottom-right: Threats
box(sl, Inches(6.75), Inches(4.05), Inches(6.1), Inches(2.8), fill=quad_col, line=RGBColor(0xFF,0x44,0x44))
txbox(sl, "⚡  Threats", Inches(6.9), Inches(4.1), Inches(5.8), Inches(0.45),
      font_size=15, bold=True, colour=RGBColor(0xFF,0x44,0x44))
bullet_list(sl,
    ["Competition from established horror titles",
     "Motion sickness risk in FPS underwater view",
     "NavMesh limitations in complex 3D mazes",
     "Time constraints for FYP deadline"],
    Inches(6.9), Inches(4.55), Inches(5.8), Inches(2.1), font_size=13, icon="▪")


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 25 — Future Scope
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
section_header(sl, "Future Scope & Enhancements")

future = [
    ("🔊 Full Audio System",
     "Ambient underwater soundscape, dynamic music that intensifies near enemies, and spatial audio for enemy footsteps."),
    ("🌐 Backend Authentication",
     "Replace PlayerPrefs with Firebase Auth or Supabase for real user accounts, cloud save, and password hashing."),
    ("🗺️ Level Expansion",
     "Design 5+ additional levels with new terrain types (caves, shipwrecks), additional enemy variants, and boss encounters."),
    ("📱 Mobile Port",
     "Port to Android using Unity's mobile build pipeline with touch-based joystick controls and accelerometer sonar."),
    ("🏆 Global Leaderboard",
     "Online score submission with real-time leaderboard using Firebase Firestore or PlayFab."),
    ("🤝 Co-op Mode",
     "Two-player cooperative survival using Unity Netcode for GameObjects — share oxygen, revive each other."),
]
y = Inches(1.05)
for icon_title, desc in future:
    box(sl, Inches(0.4), y, Inches(12.5), Inches(0.92), fill=DARK_PANEL, line=TEAL)
    txbox(sl, icon_title, Inches(0.6), y + Inches(0.03), Inches(12.1), Inches(0.38),
          font_size=14, bold=True, colour=TEAL)
    txbox(sl, desc, Inches(0.6), y + Inches(0.42), Inches(12.1), Inches(0.45),
          font_size=13, colour=WHITE)
    y += Inches(1.04)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 26 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)

box(sl, Inches(0.5), Inches(1.8), Inches(12.33), Inches(3.2), fill=DARK_PANEL, line=TEAL)
txbox(sl, "Thank You!", Inches(0.5), Inches(1.9), Inches(12.33), Inches(1.2),
      font_size=60, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)

quote = (
    '"The best games are not those with the highest budgets,\n'
    'but those with the deepest ideas."'
)
txbox(sl, quote, Inches(1.5), Inches(3.15), Inches(10.33), Inches(1.2),
      font_size=18, colour=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "— Anonymous", Inches(1.5), Inches(4.3), Inches(10.33), Inches(0.5),
      font_size=14, bold=True, colour=CYAN_LIGHT, align=PP_ALIGN.CENTER)

accent_bar(sl, top=Inches(5.2))
txbox(sl, "DarkSea Prototype  |  FYP 2022–2026",
      Inches(0.5), Inches(5.35), Inches(12.33), Inches(0.4),
      font_size=14, colour=CYAN_LIGHT, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\ali41\OneDrive\Desktop\backup game\DarkSea_Prototype\DarkSea_FYP_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
